import io
import json
import html

import streamlit as st
import anthropic
import pdfplumber
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        lines.append(para.text)
    return "\n".join(lines)


def extract_info(text: str, client: anthropic.Anthropic) -> dict:
    response = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=512,
        messages=[{
            "role": "user",
            "content": (
                "从以下文档中提取企业信息，以JSON格式返回，每项内容尽量精简（30字以内），"
                "找不到的字段填\"—\"。\n\n"
                "返回格式：\n"
                "{\"company_name\":\"企业名称\",\"core_team\":\"核心团队\","
                "\"core_business\":\"核心业务\",\"business_progress\":\"业务进展\","
                "\"track_overview\":\"赛道概况\"}\n\n"
                f"文档内容：\n{text[:6000]}\n\n仅返回JSON，不要其他文字。"
            )
        }]
    )
    raw = response.content[0].text.strip()
    # 去除可能的 markdown 代码块
    if "```" in raw:
        raw = raw.split("```")[1].lstrip("json").strip()
        raw = raw.split("```")[0].strip()
    return json.loads(raw)


def render_table(results: list) -> str:
    th = "background:#f0f2f6;padding:10px 14px;border:1px solid #ddd;text-align:left;font-weight:600"
    td_name = "padding:10px 14px;border:1px solid #ddd;vertical-align:top;font-weight:600;white-space:nowrap"
    td_info = "padding:10px 14px;border:1px solid #ddd;vertical-align:top;line-height:2"

    rows = ""
    for r in results:
        info_html = (
            f"<b>核心团队：</b>{html.escape(r.get('core_team', '—'))}<br>"
            f"<b>核心业务：</b>{html.escape(r.get('core_business', '—'))}<br>"
            f"<b>业务进展：</b>{html.escape(r.get('business_progress', '—'))}<br>"
            f"<b>赛道概况：</b>{html.escape(r.get('track_overview', '—'))}"
        )
        rows += (
            f"<tr>"
            f"<td style='{td_name}'>{html.escape(r.get('company_name', '未识别'))}</td>"
            f"<td style='{td_info}'>{info_html}</td>"
            f"</tr>"
        )

    return (
        f"<table style='width:100%;border-collapse:collapse;margin-top:12px'>"
        f"<thead><tr>"
        f"<th style='{th}'>企业名称</th>"
        f"<th style='{th}'>核心信息</th>"
        f"</tr></thead>"
        f"<tbody>{rows}</tbody>"
        f"</table>"
    )


def main():
    st.set_page_config(page_title="企业信息提取", page_icon="📊", layout="wide")
    st.title("📊 企业信息提取")
    st.caption("上传 PDF 或 PPTX 文件，自动提取企业名称、核心团队、核心业务、业务进展、赛道概况")

    with st.sidebar:
        st.header("配置")
        api_key = st.text_input("Anthropic API Key", type="password")

    uploaded_files = st.file_uploader(
        "上传文件（支持 PDF 和 PPTX，可多选）",
        type=["pdf", "pptx"],
        accept_multiple_files=True,
    )

    if not uploaded_files:
        return

    if not api_key:
        st.warning("请在左侧侧栏输入 Anthropic API Key")
        return

    if st.button("开始提取", type="primary"):
        client = anthropic.Anthropic(api_key=api_key)
        results = []

        for f in uploaded_files:
            with st.spinner(f"正在处理：{f.name}"):
                try:
                    data = f.read()
                    if f.name.lower().endswith(".pdf"):
                        text = extract_text_from_pdf(data)
                    else:
                        text = extract_text_from_pptx(data)

                    if not text.strip():
                        st.warning(f"{f.name}：未能提取到文本内容")
                        continue

                    info = extract_info(text, client)
                    results.append(info)

                except json.JSONDecodeError:
                    st.error(f"{f.name}：解析 AI 返回结果失败，请重试")
                except Exception as e:
                    st.error(f"{f.name}：{e}")

        if results:
            st.success(f"共处理 {len(results)} 个文件")
            st.markdown(render_table(results), unsafe_allow_html=True)


if __name__ == "__main__":
    main()
